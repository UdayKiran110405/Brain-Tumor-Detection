from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "HYBRID CNN-RF AND CNN-XGBOOST FRAMEWORK\nFOR REAL-TIME FOREST FIRE DETECTION"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(32)
    
    subtitle.text = "N.SUDHA RANI, V.GAYATHRI, G.VARA PRASAD, P.HEMA SRI, D. NARENDAR\n" \
                   "Sai Spurthi Institute Of Technology, Khammam, Telangana, India"
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Slide 2: Abstract
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "ABSTRACT"
    content.text = ("• Primary goal: Detect forest fires as early as possible\n"
                    "• Uses camera images and movies for detection\n"
                    "• Focuses on picture recognition algorithms\n"
                    "• Process: Background subtraction → Color segmentation → CNN classification\n"
                    "• Automatically labels fires and alerts forest service\n"
                    "• Achieves 98% accuracy in fire detection")
    
    # Slide 3: Introduction (Part 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "INTRODUCTION"
    content.text = ("• Forest fires among most catastrophic natural occurrences\n"
                    "• Climate change increasing frequency/destruction\n"
                    "• 36% of India's forests classified as flammable\n"
                    "• Need for early detection to minimize losses\n"
                    "• Computer vision-based systems becoming prevalent\n"
                    "• Cameras now widely used for fire detection")
    
    # Slide 4: Introduction (Part 2)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "INTRODUCTION (CONT.)"
    content.text = ("• Project reduces manual intervention using CNN\n"
                    "• CNNs autonomously learn/extract visual information\n"
                    "• 98% accuracy in fire photo identification\n"
                    "• System uses color segmentation + background removal\n"
                    "• Automatically detects fires and sends alerts\n"
                    "• Effective in various environmental conditions")
    
    # Slide 5: Literature Survey (Overview)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "LITERATURE SURVEY"
    content.text = ("• Review of 15 recent papers (2020-2024)\n"
                    "• Common approaches: CNN-RF and CNN-XGBoost hybrids\n"
                    "• Data sources: Satellite, drone, thermal images\n"
                    "• Key findings:\n"
                    "  - Hybrid models improve accuracy\n"
                    "  - Effective in remote areas\n"
                    "  - Reduce false positives\n"
                    "  - Faster detection than traditional methods")
    
    # Slide 6: Literature Survey (Key Papers 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "KEY RESEARCH PAPERS (1)"
    content.text = ("1. Zhang & Liu (2024): CNN-RF-XGBoost hybrid\n"
                    "   - Satellite images + ML\n"
                    "   - Enhanced speed/accuracy\n\n"
                    "2. He & Li (2023): CNN-RF/CNN-XGBoost\n"
                    "   - Remote sensing data\n"
                    "   - Effective in remote areas\n\n"
                    "3. Smith & Chang (2023): Real-time CNN-RF\n"
                    "   - Drone/satellite images\n"
                    "   - Large-scale monitoring")
    
    # Slide 7: Literature Survey (Key Papers 2)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "KEY RESEARCH PAPERS (2)"
    content.text = ("4. Lee & Kim (2022): CNN vs XGBoost\n"
                    "   - Better accuracy in data-scarce areas\n\n"
                    "5. Yu & Zhang (2022): IoT + CNN-XGBoost\n"
                    "   - Low-cost system\n\n"
                    "6. Wang & Zhou (2022): Satellite imagery\n"
                    "   - Ensemble learning benefits\n\n"
                    "7. Wu & Chen (2021): CNN-XGBoost best\n"
                    "   - Fast detection for large forests")
    
    # Slide 8: Existing System
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "EXISTING SYSTEM"
    content.text = ("• Current methods:\n"
                    "  - Smoke detectors (ineffective for forests)\n"
                    "  - SVM (68% accuracy)\n"
                    "• Limitations:\n"
                    "  - Time-consuming\n"
                    "  - Unreliable results\n"
                    "  - Large datasets problematic\n"
                    "  - Sensitive to noise/parameters\n"
                    "  - No classification probabilities")
    
    # Slide 9: Proposed System
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "PROPOSED SYSTEM"
    content.text = ("• CNN-based fire detection\n"
                    "• Advantages:\n"
                    "  - 98% accuracy\n"
                    "  - No sensors/human intervention\n"
                    "  - Lower processing requirements\n"
                    "  - No need for large training sets\n"
                    "  - Excellent image recognition\n"
                    "• Automatically alerts forest department")
    
    # Slide 10: Algorithm Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "ALGORITHM OVERVIEW"
    content.text = ("1. Dataset collection\n"
                    "2. Data preprocessing\n"
                    "3. Dataset partitioning\n"
                    "4. CNN architecture design\n"
                    "5. Model compilation\n"
                    "6. Model training\n"
                    "7. Hyperparameter tuning\n"
                    "8. Model testing\n"
                    "9. Threshold setting\n"
                    "10. Model deployment\n"
                    "11. Continuous improvement")
    
    # Slide 11: CNN Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CNN ARCHITECTURE"
    content.text = ("• Three main components:\n"
                    "  1. Convolution layers\n"
                    "  2. Pooling layers\n"
                    "  3. Fully connected layers\n"
                    "• Activation functions\n"
                    "• Feature extraction:\n"
                    "  - Color segmentation\n"
                    "  - Background subtraction\n"
                    "  - Fire boundary detection")
    
    # Slide 12: Flow Chart
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "SYSTEM FLOW CHART"
    # Note: You would add the flow chart image here
    # left = top = Inches(1)
    # pic = slide.shapes.add_picture("flow_chart.png", left, top, height=Inches(5.5))
    
    # Slide 13: Results (Before Detection)
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "RESULTS: BEFORE FIRE DETECTION"
    # Note: Add before detection image here
    # pic = slide.shapes.add_picture("before_detection.png", Inches(1), Inches(1.5), height=Inches(4.5))
    
    # Slide 14: Results (After Detection)
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "RESULTS: AFTER FIRE DETECTION"
    # Note: Add after detection image here
    # pic = slide.shapes.add_picture("after_detection.png", Inches(1), Inches(1.5), height=Inches(4.5))
    
    # Slide 15: Results Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "RESULTS ANALYSIS"
    content.text = ("• Implemented Automatic Fire Detection system\n"
                    "• CNN algorithm with color segmentation\n"
                    "• Key achievements:\n"
                    "  - High accuracy rate\n"
                    "  - Effective background removal\n"
                    "  - Precise fire location identification\n"
                    "  - Automatic notification system\n"
                    "• Outperforms traditional methods")
    
    # Slide 16: Advantages
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "ADVANTAGES"
    content.text = ("• No human intervention needed\n"
                    "• No physical sensors required\n"
                    "• Works in various environments\n"
                    "• High accuracy (98%)\n"
                    "• Fast processing\n"
                    "• Continuous learning capability\n"
                    "• Cost-effective solution\n"
                    "• Scalable for large areas")
    
    # Slide 17: Conclusion (Part 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CONCLUSION"
    content.text = ("• Enhanced CNN with object recognition\n"
                    "• Achieves project objectives:\n"
                    "  - Locates forest fires\n"
                    "  - Minimizes destruction costs\n"
                    "• Early notification enables quick response\n"
                    "• Durable and versatile system\n"
                    "• Differentiates burning/non-burning areas\n"
                    "• Learns complex visual patterns effectively")
    
    # Slide 18: Conclusion (Part 2)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CONCLUSION (CONT.)"
    content.text = ("• Compelling approach to early warning\n"
                    "• Reduces catastrophic impacts\n"
                    "• Can be integrated with existing systems\n"
                    "• Handles increasing complexity\n"
                    "• Future work:\n"
                    "  - Expand to other disaster detection\n"
                    "  - Improve real-time processing\n"
                    "  - Enhance for low-visibility conditions")
    
    # Slide 19: References (Part 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "REFERENCES (1-8)"
    content.text = ("1. Zhang & Liu (2024). IEEE Access, 12, 32014-32026\n"
                    "2. He & Li (2023). Springer Nature, 5(2), 134-145\n"
                    "3. Smith & Chang (2023). J. AI & ML, 17(3), 243-259\n"
                    "4. Lee & Kim (2022). Env. Data Sci., 8(5), 1101-1116\n"
                    "5. Yu & Zhang (2022). IEEE Trans., 68(8), 7569-7577\n"
                    "6. Wang & Zhou (2022). Remote Sensing, 16, 57-67\n"
                    "7. Li & Yu (2021). J. Fire Sci., 38(3), 247-263\n"
                    "8. Zhang & Zhao (2021). IEEE Access, 8, 34112-34122")
    
    # Slide 20: References (Part 2)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "REFERENCES (9-15)"
    content.text = ("9. Wu & Chen (2021). Env. Sci. & Tech., 55(13), 8892-8903\n"
                    "10. Shah & Mehta (2021). J. Fire Tech., 57(1), 23-36\n"
                    "11. Lee & Kim (2020). Sensors, 21(4), 1267\n"
                    "12. Chao & Sun (2020). Comp. & Env., 88, 101677\n"
                    "13. Xie & Tang (2020). Env. Monitoring, 196(4), 258-274\n"
                    "14. Zhou & Zhang (2020). J. ML Research, 22(123), 1-22\n"
                    "15. Liu & Chen (2020). Env. Intelligence, 6(2), 57-67")
    
    # Save presentation
    prs.save('Forest_Fire_Detection_Presentation.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()